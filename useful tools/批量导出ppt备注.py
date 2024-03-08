# -*- coding: utf-8 -*-
"""
Created on Fri Mar  8 10:35:49 2024

@author: DELL
"""

from pptx import Presentation
outfilepath = ""
# 打开要提取备注的 PPT 文件
prs = Presentation(outfilepath)
 
# 遍历 PPT 中的每一页
with open("ppt.txt",'w',encoding="utf-8")as f:
    for slide in prs.slides:
    # 遍历每一页中的每一个备注
        for note in slide.notes_slide.notes_text_frame.paragraphs:
        # 打印备注文本
            print(note.text)
            f.write(note.text+"\n")